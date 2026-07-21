"""
core/office_sync.py
-------------------
Two-way editing of Office files that live in OneDrive/SharePoint, so the
Documentation tab can show what changed in Word/Excel/PowerPoint *and* let you
change it back.

How it works
------------
Microsoft Graph gives us the raw file bytes and takes them back on PUT. The
in-between step is format-aware parsing:

    Word        python-docx   → a list of paragraphs (text + style)
    Excel       openpyxl      → sheets of cell rows
    PowerPoint  python-pptx   → slides of text frames

Edits come back keyed by the same indices we handed out, get applied to the
original bytes (so images, themes and formatting survive — we only touch the
text/values), and the whole file is uploaded again.

Figma is read-only by design: its REST API has no write surface, so we surface
document structure and thumbnails instead of pretending you can edit here.

Every parse is defensive: a corrupt or password-locked file returns a clean
error rather than a traceback.
"""

from __future__ import annotations

import io
from typing import Any

import requests

from core.connectors import access_token

TIMEOUT = 45
GRAPH = "https://graph.microsoft.com/v1.0"
MAX_BYTES = 25_000_000       # 25 MB — beyond this the round-trip isn't worth it


class OfficeError(Exception):
    """Anything the UI should see as a readable message."""


def _headers(extra: dict[str, str] | None = None) -> dict[str, str]:
    h = {"Authorization": f"Bearer {access_token('microsoft')}"}
    h.update(extra or {})
    return h


def kind_of(name: str) -> str:
    ext = ("." + name.rsplit(".", 1)[-1].lower()) if "." in name else ""
    return {
        ".docx": "word", ".doc": "word",
        ".xlsx": "excel", ".xlsm": "excel",
        ".pptx": "powerpoint",
    }.get(ext, "other")


# ── Graph transport ──────────────────────────────────────────────────────────
def meta(item_id: str) -> dict[str, Any]:
    r = requests.get(f"{GRAPH}/me/drive/items/{item_id}", headers=_headers(), timeout=TIMEOUT)
    if r.status_code == 404:
        raise OfficeError("that file no longer exists in the drive")
    r.raise_for_status()
    d = r.json()
    return {
        "id": d.get("id"),
        "name": d.get("name"),
        "kind": kind_of(d.get("name", "")),
        "url": d.get("webUrl"),
        "size": d.get("size"),
        "modified": d.get("lastModifiedDateTime"),
        "modified_by": ((d.get("lastModifiedBy") or {}).get("user") or {}).get("displayName"),
    }


def download(item_id: str) -> bytes:
    r = requests.get(
        f"{GRAPH}/me/drive/items/{item_id}/content",
        headers=_headers(), timeout=TIMEOUT, allow_redirects=True,
    )
    if r.status_code >= 400:
        raise OfficeError(f"download failed ({r.status_code})")
    if len(r.content) > MAX_BYTES:
        raise OfficeError("file is too large to edit here")
    return r.content


def upload(item_id: str, data: bytes) -> dict[str, Any]:
    r = requests.put(
        f"{GRAPH}/me/drive/items/{item_id}/content",
        headers=_headers({"Content-Type": "application/octet-stream"}),
        data=data, timeout=TIMEOUT,
    )
    if r.status_code >= 400:
        raise OfficeError(f"save failed ({r.status_code}): {r.text[:200]}")
    d = r.json()
    return {"ok": True, "modified": d.get("lastModifiedDateTime"), "size": d.get("size")}


# ── Word ─────────────────────────────────────────────────────────────────────
def _read_docx(data: bytes) -> dict[str, Any]:
    try:
        import docx
    except ImportError as e:
        raise OfficeError("python-docx not installed — pip install python-docx") from e
    try:
        doc = docx.Document(io.BytesIO(data))
    except Exception as e:  # noqa: BLE001
        raise OfficeError(f"could not open document: {e}") from e

    paragraphs = [
        {"i": i, "text": p.text, "style": p.style.name if p.style else "Normal"}
        for i, p in enumerate(doc.paragraphs)
    ]
    tables = [
        {
            "i": ti,
            "rows": [[c.text for c in row.cells] for row in t.rows],
        }
        for ti, t in enumerate(doc.tables)
    ]
    return {"kind": "word", "paragraphs": paragraphs, "tables": tables}


def _write_docx(data: bytes, edits: dict[str, Any]) -> bytes:
    import docx
    doc = docx.Document(io.BytesIO(data))
    for key, text in (edits.get("paragraphs") or {}).items():
        i = int(key)
        if 0 <= i < len(doc.paragraphs):
            p = doc.paragraphs[i]
            if p.runs:
                # keep the first run's formatting, drop the rest
                p.runs[0].text = text
                for extra in p.runs[1:]:
                    extra.text = ""
            else:
                p.add_run(text)
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


# ── Excel ────────────────────────────────────────────────────────────────────
MAX_ROWS, MAX_COLS = 200, 40


def _read_xlsx(data: bytes) -> dict[str, Any]:
    try:
        import openpyxl
    except ImportError as e:
        raise OfficeError("openpyxl not installed — pip install openpyxl") from e
    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=False)
    except Exception as e:  # noqa: BLE001
        raise OfficeError(f"could not open workbook: {e}") from e

    sheets = []
    for ws in wb.worksheets:
        rows = [
            [("" if c.value is None else str(c.value)) for c in row]
            for row in ws.iter_rows(max_row=MAX_ROWS, max_col=MAX_COLS)
        ]

        # openpyxl pads to the requested window; trim back to what's actually
        # used so the UI doesn't render a wall of empty cells. Always leave a
        # spare row/column so there's somewhere to type.
        used_cols = max((len(r) - next((i for i, v in enumerate(reversed(r)) if v), len(r))
                         for r in rows), default=0)
        used_rows = len(rows) - next(
            (i for i, r in enumerate(reversed(rows)) if any(r)), len(rows)
        )
        cols = min(max(used_cols + 1, 3), MAX_COLS)
        rows = [r[:cols] for r in rows[: min(max(used_rows + 1, 3), MAX_ROWS)]]

        sheets.append({
            "name": ws.title,
            "rows": rows,
            "truncated": ws.max_row > MAX_ROWS or ws.max_column > MAX_COLS,
        })
    return {"kind": "excel", "sheets": sheets}


def _write_xlsx(data: bytes, edits: dict[str, Any]) -> bytes:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data))
    # edits.cells = { "Sheet1": { "2,3": "value" } }  (0-based row,col)
    for sheet_name, cells in (edits.get("cells") or {}).items():
        if sheet_name not in wb.sheetnames:
            continue
        ws = wb[sheet_name]
        for key, value in cells.items():
            r, c = (int(x) for x in key.split(","))
            cell = ws.cell(row=r + 1, column=c + 1)
            # numbers stay numbers; formulas keep their leading =
            if isinstance(value, str) and value and not value.startswith("="):
                try:
                    cell.value = int(value) if value.lstrip("-").isdigit() else float(value)
                except ValueError:
                    cell.value = value
            else:
                cell.value = value
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


# ── PowerPoint ───────────────────────────────────────────────────────────────
def _read_pptx(data: bytes) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise OfficeError("python-pptx not installed — pip install python-pptx") from e
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:  # noqa: BLE001
        raise OfficeError(f"could not open deck: {e}") from e

    slides = []
    for si, slide in enumerate(prs.slides):
        shapes = []
        for shi, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            shapes.append({
                "i": shi,
                "name": shape.name,
                "text": shape.text_frame.text,
                "placeholder": bool(getattr(shape, "is_placeholder", False)),
            })
        title = next((s["text"] for s in shapes if s["text"]), f"Slide {si + 1}")
        slides.append({"i": si, "title": title, "shapes": shapes})
    return {"kind": "powerpoint", "slides": slides}


def _write_pptx(data: bytes, edits: dict[str, Any]) -> bytes:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides = list(prs.slides)
    # edits.slides = { "0": { "3": "new text" } }  slide index -> shape index
    for s_key, shapes in (edits.get("slides") or {}).items():
        si = int(s_key)
        if not (0 <= si < len(slides)):
            continue
        shape_list = list(slides[si].shapes)
        for sh_key, text in shapes.items():
            shi = int(sh_key)
            if not (0 <= shi < len(shape_list)):
                continue
            shape = shape_list[shi]
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            para = tf.paragraphs[0]
            if para.runs:
                para.runs[0].text = text
                for extra in para.runs[1:]:
                    extra.text = ""
            else:
                para.add_run().text = text
            # clear any trailing paragraphs so old text doesn't linger
            for extra_para in tf.paragraphs[1:]:
                for run in extra_para.runs:
                    run.text = ""
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ── public: open / save ──────────────────────────────────────────────────────
_READERS = {"word": _read_docx, "excel": _read_xlsx, "powerpoint": _read_pptx}
_WRITERS = {"word": _write_docx, "excel": _write_xlsx, "powerpoint": _write_pptx}


def open_document(item_id: str) -> dict[str, Any]:
    info = meta(item_id)
    reader = _READERS.get(info["kind"])
    if not reader:
        raise OfficeError(f"{info['name']} isn't a Word, Excel or PowerPoint file")
    content = reader(download(item_id))
    return {**info, "content": content}


def save_document(item_id: str, edits: dict[str, Any]) -> dict[str, Any]:
    info = meta(item_id)
    writer = _WRITERS.get(info["kind"])
    if not writer:
        raise OfficeError("that file type can't be edited here")
    new_bytes = writer(download(item_id), edits or {})
    result = upload(item_id, new_bytes)
    return {**result, "id": item_id, "name": info["name"]}


# ── Figma (read-only) ────────────────────────────────────────────────────────
def figma_file(key: str) -> dict[str, Any]:
    """Top-level pages and frames of a Figma file, plus thumbnails."""
    tok = access_token("figma")
    r = requests.get(
        f"https://api.figma.com/v1/files/{key}?depth=2",
        headers={"Authorization": f"Bearer {tok}"}, timeout=TIMEOUT,
    )
    if r.status_code >= 400:
        raise OfficeError(f"Figma said {r.status_code}: {r.text[:160]}")
    d = r.json()

    pages = []
    frame_ids: list[str] = []
    for page in (d.get("document", {}).get("children") or []):
        frames = []
        for node in (page.get("children") or [])[:24]:
            frames.append({"id": node.get("id"), "name": node.get("name"), "type": node.get("type")})
            frame_ids.append(node.get("id"))
        pages.append({"id": page.get("id"), "name": page.get("name"), "frames": frames})

    thumbs: dict[str, str] = {}
    if frame_ids:
        try:
            ir = requests.get(
                f"https://api.figma.com/v1/images/{key}",
                headers={"Authorization": f"Bearer {tok}"},
                params={"ids": ",".join(frame_ids[:24]), "format": "png", "scale": "1"},
                timeout=TIMEOUT,
            )
            if ir.ok:
                thumbs = {k: v for k, v in (ir.json().get("images") or {}).items() if v}
        except requests.RequestException:
            pass

    return {
        "kind": "figma",
        "name": d.get("name"),
        "modified": d.get("lastModified"),
        "version": d.get("version"),
        "url": f"https://www.figma.com/file/{key}",
        "pages": pages,
        "thumbnails": thumbs,
        "readonly": True,
    }
